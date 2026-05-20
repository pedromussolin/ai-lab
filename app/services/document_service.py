"""Document service - upload, parse, chunk, embed, index."""

import io
import uuid
from pathlib import Path
from typing import Any

from app.core.config import settings
from app.core.exceptions import DocumentProcessingError
from app.models.document import Document
from app.rag.chunker import RecursiveChunker
from app.rag.indexer import AzureSearchIndexer, LocalVectorIndexer
from app.repositories.document import DocumentRepository
from app.services.embedding_service import EmbeddingService
from app.services.rag_service import RAGService


class DocumentParser:
    """Parse various document formats to plain text."""

    async def parse(self, content: bytes, filename: str, content_type: str) -> str:
        ext = Path(filename).suffix.lower()
        if ext == ".pdf" or "pdf" in content_type:
            return self._parse_pdf(content)
        elif ext in (".docx", ".doc"):
            return self._parse_docx(content)
        elif ext in (".txt", ".md", ".py", ".js", ".ts", ".sql", ".yaml", ".yml", ".json"):
            return content.decode("utf-8", errors="replace")
        elif ext in (".xlsx", ".xls"):
            return self._parse_xlsx(content)
        elif ext == ".pptx":
            return self._parse_pptx(content)
        else:
            try:
                return content.decode("utf-8", errors="replace")
            except Exception as e:
                raise DocumentProcessingError(filename, f"Unsupported format: {ext}") from e

    def _parse_pdf(self, content: bytes) -> str:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)

    def _parse_docx(self, content: bytes) -> str:
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(content))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _parse_xlsx(self, content: bytes) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                rows.append("\t".join(str(c or "") for c in row))
        return "\n".join(rows)

    def _parse_pptx(self, content: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            slides.append(f"## Slide {i + 1}\n" + "\n".join(texts))
        return "\n\n".join(slides)


class DocumentService:
    def __init__(
        self,
        doc_repo: DocumentRepository,
        embedding_svc: EmbeddingService,
        rag_svc: RAGService,
    ) -> None:
        self._repo = doc_repo
        self._embedding_svc = embedding_svc
        self._rag_svc = rag_svc
        self._parser = DocumentParser()
        self._chunker = RecursiveChunker(chunk_size=1000, chunk_overlap=200)

    async def process_document(
        self,
        document_id: uuid.UUID,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        indexing_strategy: str = "local",
    ) -> Document:
        doc = await self._repo.get_by_id(document_id)
        if not doc:
            raise ValueError(f"Document {document_id} not found")

        # Update status
        doc.status = "processing"
        await self._repo._db.flush()

        try:
            # Read file
            file_path = Path(settings.upload_dir) / str(document_id)
            content = file_path.read_bytes()

            # Parse
            text = await self._parser.parse(content, doc.filename, doc.content_type)

            # Chunk
            chunker = RecursiveChunker(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
            chunks = chunker.chunk(text, metadata={"filename": doc.filename})

            # Embed
            texts = [c.content for c in chunks]
            embeddings = await self._embedding_svc.embed_batch(texts)

            # Index
            pipeline = await self._rag_svc.get_pipeline()
            await pipeline._local_retriever._db  # ensure db is available
            local_indexer = LocalVectorIndexer(pipeline._local_retriever._db)
            await local_indexer.index_chunks(
                chunks=chunks,
                embeddings=embeddings,
                document_id=str(document_id),
            )

            if indexing_strategy in ("azure", "hybrid") and settings.azure_search_endpoint:
                azure_indexer = AzureSearchIndexer()
                await azure_indexer.ensure_index()
                await azure_indexer.index_chunks(
                    chunks=chunks,
                    embeddings=embeddings,
                    document_id=str(document_id),
                    source=doc.filename,
                )

            doc.status = "ready"
            doc.chunk_count = len(chunks)

        except Exception as e:
            doc.status = "failed"
            doc.error_message = str(e)

        await self._repo._db.flush()
        return doc
